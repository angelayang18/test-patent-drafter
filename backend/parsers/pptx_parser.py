"""Extract text from PPTX files using python-pptx."""

import re
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape


def _clean_extracted_text(text: str) -> str:
    """Normalize whitespace: trim lines and collapse runs of blank lines."""
    lines: list[str] = []
    prev_blank = False
    for raw_line in text.splitlines():
        line = re.sub(r"[ \t]+", " ", raw_line).strip()
        if not line:
            if not prev_blank:
                lines.append("")
            prev_blank = True
        else:
            lines.append(line)
            prev_blank = False
    while lines and lines[0] == "":
        lines.pop(0)
    while lines and lines[-1] == "":
        lines.pop()
    return "\n".join(lines)


def _slide_title(slide) -> str:
    """Return the slide title text, if a title placeholder exists."""
    if slide.shapes.title is None:
        return ""
    return slide.shapes.title.text.strip()


def _is_title_placeholder(shape: BaseShape) -> bool:
    """Return True if the shape is a slide title placeholder."""
    if not getattr(shape, "is_placeholder", False):
        return False
    placeholder_type = shape.placeholder_format.type
    return placeholder_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)


def _shape_text_lines(shape: BaseShape) -> list[str]:
    """Collect non-empty text from a shape, including grouped shapes."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP and isinstance(shape, GroupShape):
        lines: list[str] = []
        for child in shape.shapes:
            lines.extend(_shape_text_lines(child))
        return lines

    if _is_title_placeholder(shape) or not hasattr(shape, "text"):
        return []

    text = shape.text.strip()
    return [text] if text else []


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract all text from a PowerPoint file provided as raw bytes.

    Each slide is introduced with a header of the form ``=== Slide N: Title ===``.
    Body text from shapes on the slide follows the header (title text is not repeated).

    Args:
        file_bytes: Raw PPTX file contents.

    Returns:
        Extracted slide text as a single cleaned string.

    Raises:
        ValueError: If ``file_bytes`` is empty.
        RuntimeError: If the presentation cannot be opened or text extraction fails.
    """
    if not file_bytes:
        raise ValueError("PPTX input is empty; cannot extract text from zero bytes.")

    try:
        presentation = Presentation(BytesIO(file_bytes))
        slide_sections: list[str] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = _slide_title(slide)
            header = (
                f"=== Slide {slide_number}: {title} ==="
                if title
                else f"=== Slide {slide_number} ==="
            )

            body_lines: list[str] = []
            for shape in slide.shapes:
                body_lines.extend(_shape_text_lines(shape))

            section_parts = [header]
            if body_lines:
                section_parts.append("\n".join(body_lines))
            slide_sections.append("\n".join(section_parts))

        return _clean_extracted_text("\n\n".join(slide_sections))
    except ValueError:
        raise
    except Exception as exc:
        raise RuntimeError(f"Failed to extract text from PPTX: {exc}") from exc
